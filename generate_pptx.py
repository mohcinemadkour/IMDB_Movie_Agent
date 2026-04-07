"""
Generate app_presentation.pptx from the content of app.md.
Run with:  .venv\Scripts\python generate_pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette ────────────────────────────────────────────────────────────
DARK_BG   = RGBColor(0x1C, 0x1C, 0x2E)   # deep navy
ACCENT    = RGBColor(0xFF, 0x4B, 0x4B)   # IMDB-ish red
ACCENT2   = RGBColor(0xF5, 0xC5, 0x18)   # IMDB gold
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY= RGBColor(0xCC, 0xCC, 0xDD)
MID_GREY  = RGBColor(0x88, 0x88, 0x99)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]   # totally blank


# ── Helper utilities ──────────────────────────────────────────────────────────

def add_rect(slide, left, top, width, height, fill_rgb, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text,
                font_size=18, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox


def add_bullet_slide(title_text, bullets, subtitle=None):
    slide = prs.slides.add_slide(blank_layout)

    # Background
    add_rect(slide, 0, 0, 13.33, 7.5, DARK_BG)
    # Left accent bar
    add_rect(slide, 0, 0, 0.07, 7.5, ACCENT)
    # Title bar
    add_rect(slide, 0.07, 0, 13.26, 1.1, RGBColor(0x14, 0x14, 0x24))

    # Title
    add_textbox(slide, 0.25, 0.1, 12.8, 0.9, title_text,
                font_size=28, bold=True, color=ACCENT2, align=PP_ALIGN.LEFT)

    if subtitle:
        add_textbox(slide, 0.25, 1.15, 12.8, 0.4, subtitle,
                    font_size=14, color=MID_GREY, italic=True)

    # Bullets
    y = 1.5 if not subtitle else 1.9
    for bullet in bullets:
        # bullet dot
        add_rect(slide, 0.35, y + 0.08, 0.08, 0.08, ACCENT)
        add_textbox(slide, 0.55, y, 12.2, 0.5, bullet,
                    font_size=16, color=LIGHT_GREY)
        y += 0.52

    return slide


def add_table_slide(title_text, headers, rows, subtitle=None):
    slide = prs.slides.add_slide(blank_layout)

    add_rect(slide, 0, 0, 13.33, 7.5, DARK_BG)
    add_rect(slide, 0, 0, 0.07, 7.5, ACCENT)
    add_rect(slide, 0.07, 0, 13.26, 1.1, RGBColor(0x14, 0x14, 0x24))

    add_textbox(slide, 0.25, 0.1, 12.8, 0.9, title_text,
                font_size=28, bold=True, color=ACCENT2)

    if subtitle:
        add_textbox(slide, 0.25, 1.15, 12.8, 0.4, subtitle,
                    font_size=14, color=MID_GREY, italic=True)

    top_offset = 1.6 if not subtitle else 2.0
    col_count  = len(headers)
    col_width  = 12.5 / col_count

    # Header row
    for i, h in enumerate(headers):
        add_rect(slide, 0.25 + i * col_width, top_offset,
                 col_width - 0.05, 0.42, ACCENT)
        add_textbox(slide, 0.3 + i * col_width, top_offset + 0.04,
                    col_width - 0.1, 0.38, h,
                    font_size=14, bold=True, color=WHITE)

    # Data rows
    for r_idx, row in enumerate(rows):
        row_bg = RGBColor(0x24, 0x24, 0x38) if r_idx % 2 == 0 else RGBColor(0x1C, 0x1C, 0x2E)
        y = top_offset + 0.44 + r_idx * 0.44
        add_rect(slide, 0.25, y, 12.5, 0.42, row_bg)
        for c_idx, cell in enumerate(row):
            add_textbox(slide, 0.3 + c_idx * col_width, y + 0.04,
                        col_width - 0.1, 0.38, str(cell),
                        font_size=13, color=LIGHT_GREY)

    return slide


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title slide
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, DARK_BG)
add_rect(slide, 0, 3.4, 13.33, 0.06, ACCENT)

# Big clapperboard emoji substitute — coloured rectangle title block
add_rect(slide, 0.5, 1.6, 12.33, 1.6, RGBColor(0x14, 0x14, 0x24))
add_textbox(slide, 0.5, 1.65, 12.33, 1.5,
            "🎬  IMDB Movie Agent",
            font_size=44, bold=True, color=ACCENT2, align=PP_ALIGN.CENTER)

add_textbox(slide, 0.5, 3.55, 12.33, 0.7,
            "app.py — Architecture & Design Documentation",
            font_size=22, color=WHITE, align=PP_ALIGN.CENTER)

add_textbox(slide, 0.5, 4.4, 12.33, 0.5,
            "Streamlit  ·  OpenAI GPT-4o  ·  FAISS  ·  LangGraph  ·  Whisper TTS",
            font_size=14, color=MID_GREY, align=PP_ALIGN.CENTER, italic=True)

add_textbox(slide, 0.5, 6.8, 12.33, 0.4,
            "April 2026",
            font_size=12, color=MID_GREY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Overview
# ═══════════════════════════════════════════════════════════════════════════════
add_bullet_slide(
    "Overview",
    [
        "Streamlit-based conversational UI for the IMDB Top 1000 dataset",
        "Combines structured pandas queries with FAISS semantic (vector) search",
        "Powered by OpenAI GPT-4o via a LangGraph agent with 3 tools",
        "Supports voice input (Whisper STT) and voice output (TTS-1 / gTTS)",
        "Persistent chat history with example question shortcuts",
        "All API keys stored in .env — never hardcoded",
    ],
)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Startup Sequence
# ═══════════════════════════════════════════════════════════════════════════════
add_bullet_slide(
    "Startup Sequence",
    [
        "1. python-dotenv loads .env  (OPENAI_API_KEY / GOOGLE_API_KEY)",
        "2. st.set_page_config — must be the very first Streamlit call",
        "3. API key guard — st.stop() if no key is found",
        "4. Session state init — messages list + agent_executor cached once",
        "5. First run: FAISS vector index is built over the Overview column",
        "6. Subsequent runs: index loaded from data/faiss_index/ (fast)",
    ],
    subtitle="Run order matters — each step depends on the previous",
)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Module Layout
# ═══════════════════════════════════════════════════════════════════════════════
add_bullet_slide(
    "Module Layout  (app.py top-level sections)",
    [
        "_speak()              — TTS helper function",
        "Sidebar               — Settings panel + 10 example question buttons",
        "Main heading          — Title + caption bar",
        "Voice input expander  — Mic recording + Whisper transcription pipeline",
        "Chat history display  — Renders all past messages on every rerun",
        "Input handling        — chat_input / pending_input → agent → response",
    ],
)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — _speak() helper
# ═══════════════════════════════════════════════════════════════════════════════
add_table_slide(
    "_speak()  — Text-to-Speech Helper",
    ["Parameter", "Description"],
    [
        ["text",   "String to speak — 4096 char limit (OpenAI) / 500 chars (gTTS)"],
        ["engine", '"auto" (OpenAI → gTTS fallback), "openai", or "gtts"'],
        ["voice",  "OpenAI voice: nova / alloy / echo / fable / onyx / shimmer"],
        ["returns","Engine used: \"openai\" | \"gtts\" | \"error: <msg>\""],
    ],
    subtitle="st.audio(..., autoplay=True) plays the audio inline — no page reload needed",
)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Sidebar controls
# ═══════════════════════════════════════════════════════════════════════════════
add_table_slide(
    "Sidebar Controls",
    ["Control", "Purpose"],
    [
        ["Voice Output toggle",    "Enable / disable TTS after each agent response"],
        ["TTS Engine radio",       "Select Auto / OpenAI TTS / gTTS"],
        ["OpenAI Voice selectbox", "Pick from 6 voices (hidden when gTTS selected)"],
        ["Clear Chat button",      "Resets messages list + triggers st.rerun()"],
        ["Example question buttons","Sets pending_input + reruns to auto-submit"],
    ],
)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Voice Input Pipeline
# ═══════════════════════════════════════════════════════════════════════════════
add_bullet_slide(
    "Voice Input Pipeline",
    [
        "audio_recorder() captures mic bytes via audio-recorder-streamlit widget",
        "Guard: skip if bytes ≤ 1000  (header-only = no real audio)",
        "MD5 deduplication: skip if same clip already processed this session",
        "Whisper API: model=whisper-1, verbose_json, language=en, IMDB domain prompt",
        "Confidence check: avg no_speech_prob ≥ 0.6 → warn + allow re-record",
        "Per-segment 🟢/🟡/🔴 confidence expander shown to user",
        "Confirmation UI: editable text box + ✅ Send / 🗑️ Discard buttons",
    ],
    subtitle="Collapsible expander — voice input is fully optional",
)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Whisper call details
# ═══════════════════════════════════════════════════════════════════════════════
add_bullet_slide(
    "Whisper Transcription — Key Details",
    [
        'response_format="verbose_json": returns TranscriptionSegment objects (not dicts)',
        "Attribute access via getattr(seg, 'no_speech_prob', 0) — Pydantic model",
        'Domain prompt: "IMDB, movie titles, director names, genres like Horror, Comedy, Sci-Fi"',
        "no_speech_prob threshold: 0.6 — silence/noise likely above this value",
        "Low confidence → ⚠️ warning displayed + _last_audio_hash reset",
        "High confidence → store in session_state['_pending_transcription']",
    ],
    subtitle="verbose_json enables per-segment quality scoring unavailable in plain text format",
)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Session State Keys
# ═══════════════════════════════════════════════════════════════════════════════
add_table_slide(
    "Session State Keys",
    ["Key", "Type", "Purpose"],
    [
        ["messages",               "list[dict]",         "Full chat history (role + content)"],
        ["agent_executor",         "CompiledStateGraph", "Cached LangGraph agent — built once"],
        ["_last_audio_hash",       "str",                "MD5 of last processed audio (dedup)"],
        ["_pending_transcription", "str",                "Whisper output awaiting confirmation"],
        ["pending_input",          "str",                "Question ready to submit to agent"],
    ],
)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Input Handling Flow
# ═══════════════════════════════════════════════════════════════════════════════
add_bullet_slide(
    "Input Handling Flow",
    [
        "Two sources merged: st.chat_input (keyboard) + pending_input (voice/buttons)",
        "1. User message appended to history and rendered in chat bubble",
        "2. run_agent(executor, user_input, history) called inside st.spinner",
        "3. Exceptions caught — friendly error shown instead of crash",
        "4. Assistant response rendered with st.markdown",
        "5. If Voice Output on: _speak() called, status caption shown below response",
        "6. Assistant message appended to history for next turn context",
    ],
    subtitle="history passed to agent excludes the current message (already added separately)",
)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Dependencies
# ═══════════════════════════════════════════════════════════════════════════════
add_table_slide(
    "Key Dependencies",
    ["Package", "Usage"],
    [
        ["streamlit",                 "UI framework — chat, widgets, audio playback"],
        ["openai",                    "Whisper-1 STT, TTS-1 speech, GPT-4o LLM"],
        ["audio-recorder-streamlit",  "Microphone recording widget in-browser"],
        ["gtts",                      "Google TTS fallback (free, no key needed)"],
        ["langchain + langgraph",     "Agent orchestration, tool routing"],
        ["faiss-cpu",                 "Vector similarity search over Overview column"],
        ["pandas",                    "All structured data filtering and sorting"],
        ["python-dotenv",             ".env loading for API keys"],
    ],
)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Summary / Thank you
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, DARK_BG)
add_rect(slide, 0, 3.3, 13.33, 0.06, ACCENT)

add_textbox(slide, 0.5, 1.2, 12.33, 1.0,
            "Thank You",
            font_size=44, bold=True, color=ACCENT2, align=PP_ALIGN.CENTER)

add_textbox(slide, 0.5, 2.4, 12.33, 0.6,
            "IMDB Movie Agent — app.py",
            font_size=22, color=WHITE, align=PP_ALIGN.CENTER)

points = [
    "✅  Streamlit conversational UI with voice I/O",
    "✅  OpenAI GPT-4o + LangGraph agent (3 tools)",
    "✅  FAISS semantic search over 1 000 movie overviews",
    "✅  Whisper STT with confidence scoring & confirmation UI",
    "✅  OpenAI TTS + gTTS fallback with engine/voice selector",
]
for i, p in enumerate(points):
    add_textbox(slide, 1.5, 3.55 + i * 0.52, 10.5, 0.48,
                p, font_size=15, color=LIGHT_GREY, align=PP_ALIGN.LEFT)


# ── Save ──────────────────────────────────────────────────────────────────────
output = "app_presentation.pptx"
prs.save(output)
print(f"✅  Saved: {output}")
